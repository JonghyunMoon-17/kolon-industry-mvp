from __future__ import annotations

import hashlib
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import BinaryIO

from src.models import ElementType, ParsedChunk
from src.utils import split_text


PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
TEXT_EXTENSIONS = {".txt", ".md"}
SUPPORTED_EXTENSIONS = PDF_EXTENSIONS | DOCX_EXTENSIONS | PPTX_EXTENSIONS | TEXT_EXTENSIONS


def _chunk_id(source_file: str, index: int, text: str) -> str:
    digest = hashlib.sha1(f"{source_file}:{index}:{text[:80]}".encode("utf-8")).hexdigest()[:12]
    return f"chunk_{digest}"


def _make_chunks(
    *,
    source_file: str,
    doc_type: str,
    text_units: list[dict],
    default_element_type: ElementType,
) -> list[ParsedChunk]:
    chunks: list[ParsedChunk] = []
    idx = 0
    for unit in text_units:
        for text in split_text(unit.get("text", "")):
            chunks.append(
                ParsedChunk(
                    chunk_id=_chunk_id(source_file, idx, text),
                    source_file=source_file,
                    doc_type=doc_type,
                    page_or_slide=unit.get("page_or_slide"),
                    section_title=unit.get("section_title"),
                    text=text,
                    element_type=unit.get("element_type", default_element_type),
                    metadata=unit.get("metadata", {}),
                )
            )
            idx += 1
    return chunks


def parse_uploaded_file(uploaded_file) -> list[ParsedChunk]:
    suffix = Path(uploaded_file.name).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {suffix}")
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        temp_path = Path(tmp.name)
    try:
        return parse_path(temp_path, original_name=uploaded_file.name)
    finally:
        temp_path.unlink(missing_ok=True)


def parse_path(path: str | Path, original_name: str | None = None) -> list[ParsedChunk]:
    path = Path(path)
    source_file = original_name or path.name
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        return parse_text(path, source_file)
    if suffix in PDF_EXTENSIONS:
        return parse_pdf(path, source_file)
    if suffix in DOCX_EXTENSIONS:
        return parse_docx(path, source_file)
    if suffix in PPTX_EXTENSIONS:
        return parse_pptx(path, source_file)
    raise ValueError(f"Unsupported file format: {suffix}")


def parse_text(path: Path, source_file: str) -> list[ParsedChunk]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    element = ElementType.TRANSCRIPT if "transcript" in source_file.lower() or "녹취" in source_file else ElementType.TEXT
    units = [{"text": text, "metadata": {"parser": "plain_text"}}]
    return _make_chunks(source_file=source_file, doc_type=path.suffix.lower().lstrip("."), text_units=units, default_element_type=element)


def parse_pdf(path: Path, source_file: str) -> list[ParsedChunk]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PDF parsing requires PyMuPDF. Install requirements.txt.") from exc

    units: list[dict] = []
    with fitz.open(path) as doc:
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                units.append(
                    {
                        "text": text,
                        "page_or_slide": str(page_index),
                        "section_title": f"Page {page_index}",
                        "element_type": ElementType.TEXT,
                        "metadata": {"parser": "pymupdf", "page": page_index},
                    }
                )
    return _make_chunks(source_file=source_file, doc_type="pdf", text_units=units, default_element_type=ElementType.TEXT)


def parse_docx(path: Path, source_file: str) -> list[ParsedChunk]:
    try:
        from docx import Document
    except ImportError:
        return parse_docx_zip_fallback(path, source_file)

    doc = Document(path)
    units: list[dict] = []
    current_heading: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        nonlocal buffer
        text = "\n".join(buffer).strip()
        if text:
            units.append(
                {
                    "text": text,
                    "section_title": current_heading,
                    "element_type": ElementType.TEXT,
                    "metadata": {"parser": "python-docx"},
                }
            )
        buffer = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name.lower().startswith("heading"):
            flush()
            current_heading = text
        else:
            buffer.append(text)
    flush()

    for table_index, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        table_text = "\n".join(rows)
        if table_text.strip():
            units.append(
                {
                    "text": table_text,
                    "section_title": f"Table {table_index}",
                    "element_type": ElementType.TABLE,
                    "metadata": {"parser": "python-docx", "table_index": table_index},
                }
            )
    return _make_chunks(source_file=source_file, doc_type="docx", text_units=units, default_element_type=ElementType.TEXT)


def parse_docx_zip_fallback(path: Path, source_file: str) -> list[ParsedChunk]:
    """Read DOCX text without python-docx.

    DOCX files are zip archives containing WordprocessingML. This fallback keeps
    the demo usable when optional dependencies are missing.
    """

    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    units: list[dict] = []
    with zipfile.ZipFile(path) as archive:
        document_xml = archive.read("word/document.xml")
    root = ET.fromstring(document_xml)

    buffer: list[str] = []
    current_heading: str | None = None

    def flush() -> None:
        nonlocal buffer
        text = "\n".join(item for item in buffer if item.strip()).strip()
        if text:
            units.append(
                {
                    "text": text,
                    "section_title": current_heading,
                    "element_type": ElementType.TEXT,
                    "metadata": {"parser": "docx_zip_fallback"},
                }
            )
        buffer = []

    for paragraph in root.findall(".//w:p", ns):
        texts = [node.text for node in paragraph.findall(".//w:t", ns) if node.text]
        text = "".join(texts).strip()
        if not text:
            continue
        style = paragraph.find(".//w:pStyle", ns)
        style_val = style.attrib.get(f"{{{ns['w']}}}val", "") if style is not None else ""
        if "Heading" in style_val or "Title" in style_val or len(text) <= 40 and text.endswith(("분석", "현황", "시장", "산업")):
            flush()
            current_heading = text
        else:
            buffer.append(text)
    flush()
    return _make_chunks(source_file=source_file, doc_type="docx", text_units=units, default_element_type=ElementType.TEXT)


def parse_pptx(path: Path, source_file: str) -> list[ParsedChunk]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX parsing requires python-pptx. Install requirements.txt.") from exc

    prs = Presentation(path)
    units: list[dict] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        title = None
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if title is None:
                    title = shape.text.strip().splitlines()[0][:120]
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append(" | ".join(cell.text.strip() for cell in row.cells))
                if rows:
                    units.append(
                        {
                            "text": "\n".join(rows),
                            "page_or_slide": str(slide_index),
                            "section_title": title or f"Slide {slide_index} Table",
                            "element_type": ElementType.TABLE,
                            "metadata": {"parser": "python-pptx", "slide": slide_index},
                        }
                    )
        if parts:
            units.append(
                {
                    "text": "\n".join(parts),
                    "page_or_slide": str(slide_index),
                    "section_title": title or f"Slide {slide_index}",
                    "element_type": ElementType.SLIDE,
                    "metadata": {"parser": "python-pptx", "slide": slide_index},
                }
            )
    return _make_chunks(source_file=source_file, doc_type="pptx", text_units=units, default_element_type=ElementType.SLIDE)
